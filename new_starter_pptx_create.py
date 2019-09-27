from pptx import Presentation
from pptx.dml.color import RGBColor
import pandas as pd
from datetime import date, timedelta
from dateutil import parser

prs = Presentation('New_Starter_Template.pptx')

starter_list = pd.read_excel('new_starters.xlsx')
# starter_list['Date'] = starter_list['Date']
today_slide = prs.slides[0].shapes[8].text_frame

tomorrow = pd.Timestamp(date.today() + timedelta(days=1))
print(tomorrow.date().strftime('%d %B %Y'))

for index, row in starter_list.iterrows():
    if row.Date == date.today():
        p = today_slide.add_paragraph()
        p.text = f'{row.Firstname} {row.Surname} \t\t {row.Department}\n'
        p.level = 0
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 94, 184)
        print(f'{row.Firstname} {row.Surname} - {row.Department}\n')

save_name = f"Welcome {tomorrow.date().strftime('%d %B %Y')}.pptx"
prs.save(save_name)
